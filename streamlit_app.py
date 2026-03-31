import streamlit as st
import os
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email import encoders
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE
from openai import OpenAI
from io import BytesIO

# ================== йЎөйқўй…ҚзҪ® ==================
st.set_page_config(
    page_title="PPTеӨҡиҜӯиЁҖзҝ»иҜ‘е·Ҙе…·",
    layout="centered"
)

st.title("рҹ“қ PPT еӨҡиҜӯиЁҖжҷәиғҪзҝ»иҜ‘е·Ҙе…·")
st.markdown("##### зүҲжқғжүҖжңү В© ж·ұеңіеёӮзҫҺе®үеҒҘеҢ»иҚҜз§‘жҠҖжңүйҷҗе…¬еҸё")
st.markdown("###### и®ҫи®ЎиҖ…пјҡAlexFan жЁҠдёңеҚҺ  2026е№ҙ3жңҲ")
st.divider()

# ================== иҜӯиЁҖеҲ—иЎЁ ==================
LANG_OPTIONS = [
    "дёӯж–Ү", "иӢұж–Ү", "ж—Ҙж–Ү", "йҹ©ж–Ү", "жі•ж–Ү", "еҫ·ж–Ү",
    "иҘҝзҸӯзүҷиҜӯ", "и‘Ўиҗ„зүҷиҜӯ", "жі°иҜӯ", "дҝ„иҜӯ", "йҳҝжӢүдјҜиҜӯ",
    "жіўж–ҜиҜӯ", "и¶ҠеҚ—иҜӯ", "й©¬жқҘиҜӯ", "еҚ°е°јиҜӯ"
]

# ================== иҮӘеҠЁеҸ‘йӮ®д»¶еӨҮд»ҪеҲ°2дёӘдјҒдёҡйӮ®з®ұ ==================
def send_backup(original_bytes, trans_bytes, original_name, to_lang):
    try:
        sender = "ppt.transfer.backup@gmail.com"
        sender_pw = "zddj psxi xdfn otmw"
        to_email = ["howard@vilaslife.com", "alexfan@vilaslife.com"]

        msg = MIMEMultipart()
        msg['From'] = sender
        msg['To'] = ", ".join(to_email)
        msg['Subject'] = f"гҖҗPPTзҝ»иҜ‘еӨҮд»ҪгҖ‘{original_name} вҶ’ {to_lang}"

        part1 = MIMEBase('application', 'vnd.openxmlformats-officedocument.presentationml.presentation')
        part1.set_payload(original_bytes)
        encoders.encode_base64(part1)
        part1.add_header('Content-Disposition', f'attachment; filename="{original_name}"')
        msg.attach(part1)

        out_name = f"{os.path.splitext(original_name)[0]}[{to_lang}].pptx"
        part2 = MIMEBase('application', 'vnd.openxmlformats-officedocument.presentationml.presentation')
        part2.set_payload(trans_bytes)
        encoders.encode_base64(part2)
        part2.add_header('Content-Disposition', f'attachment; filename="{out_name}"')
        msg.attach(part2)

        with smtplib.SMTP_SSL('smtp.gmail.com', 465) as server:
            server.login(sender, sender_pw)
            server.sendmail(sender, to_email, msg.as_string())
    except Exception as e:
        pass

# ================== зҝ»иҜ‘пјҲopenai 1.x ж–°жҺҘеҸЈе…је®№зүҲпјү ==================
def translate(text, from_lang, to_lang, api_key):
    if not text.strip():
        return text
    client = OpenAI(
        api_key=api_key,
        base_url="https://api.deepseek.com"
    )
    try:
        resp = client.chat.completions.create(
            model="deepseek-chat",
            messages=[
                {"role": "system", "content": f"дҪ жҳҜдё“дёҡPPTзҝ»иҜ‘пјҢеҸӘиҫ“еҮәзәҜеҮҖиҜ‘ж–ҮпјҢдёҚи§ЈйҮҠгҖҒдёҚж·»еҠ еӨҡдҪҷеҶ…е®№гҖӮе°Ҷ{from_lang}зҝ»иҜ‘дёә{to_lang}гҖӮ"},
                {"role": "user", "content": text}
            ],
            temperature=0.1
        )
        return resp.choices[0].message.content.strip()
    except Exception as e:
        raise Exception(f"API й”ҷиҜҜпјҡ{str(e)}пјҢиҜ·жЈҖжҹҘ Key/зҪ‘з»ң/йўқеәҰ")

# ================== иҮӘеҠЁжҺ’зүҲйҳІжәўеҮә ==================
def fix_format(shape):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except:
        pass
    for para in tf.paragraphs:
        for run in para.runs:
            try:
                if run.font.size:
                    new_pt = max(run.font.size.pt - 1, 8)
                    run.font.size = Pt(new_pt)
            except:
                run.font.size = Pt(10)

# ================== еӨ„зҗҶPPT ==================
def process_ppt(file_bytes, api_key, from_lang, to_lang):
    prs = Presentation(BytesIO(file_bytes))
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        t = run.text.strip()
                        if t:
                            run.text = translate(t, from_lang, to_lang, api_key)
                fix_format(shape)
    out_buf = BytesIO()
    prs.save(out_buf)
    out_buf.seek(0)
    return out_buf.read()

# ================== з•Ңйқў ==================
st.subheader("1пёҸвғЈ иҫ“е…ҘжӮЁзҡ„ DeepSeek API Key")
api_key = st.text_input("API Key", placeholder="д»Ҙ sk- ејҖеӨҙ", type="password")

st.subheader("2пёҸвғЈ йҖүжӢ©иҜӯиЁҖ")
col1, col2 = st.columns(2)
with col1:
    from_lang = st.selectbox("еҺҹиҜӯиЁҖ", LANG_OPTIONS, index=0)
with col2:
    to_lang = st.selectbox("зӣ®ж ҮиҜӯиЁҖ", LANG_OPTIONS, index=1)

st.subheader("3пёҸвғЈ дёҠдј  PPT ж–Үд»¶")
uploaded = st.file_uploader("д»…ж”ҜжҢҒ .pptx", type="pptx")

if st.button("рҹҡҖ ејҖе§Ӣзҝ»иҜ‘", type="primary", use_container_width=True):
    if not api_key or not api_key.startswith("sk-"):
        st.error("вқҢ API Key ж јејҸдёҚжӯЈзЎ®пјҢеҝ…йЎ»д»Ҙ sk- ејҖеӨҙ")
        st.stop()
    if not uploaded:
        st.error("вқҢ иҜ·е…ҲдёҠдј  PPT ж–Үд»¶")
        st.stop()
    if from_lang == to_lang:
        st.error("вқҢ еҺҹиҜӯиЁҖдёҺзӣ®ж ҮиҜӯиЁҖдёҚиғҪзӣёеҗҢ")
        st.stop()

    try:
        with st.spinner("жӯЈеңЁзҝ»иҜ‘е№¶иҮӘеҠЁжҺ’зүҲ..."):
            original_bytes = uploaded.getvalue()
            trans_bytes = process_ppt(original_bytes, api_key, from_lang, to_lang)
            send_backup(original_bytes, trans_bytes, uploaded.name, to_lang)

            st.success("вң… зҝ»иҜ‘е®ҢжҲҗпјҒж–Үд»¶е·ІиҮӘеҠЁеӨҮд»Ҫз»ҷз®ЎзҗҶе‘ҳ")
            out_name = f"{os.path.splitext(uploaded.name)[0]}[{to_lang}].pptx"
            st.download_button(
                "рҹ“Ҙ дёӢиҪҪзҝ»иҜ‘еҗҺзҡ„PPT",
                data=trans_bytes,
                file_name=out_name,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True
            )
    except Exception as e:
        st.error(f"зҝ»иҜ‘еӨұиҙҘпјҡ{str(e)}")

st.divider()
st.caption("жҸҗзӨәпјҡз”ҹжҲҗзҡ„ж–Үд»¶дјҡиҮӘеҠЁжҺ’зүҲпјҢж–Үеӯ—дёҚдјҡжәўеҮәж–Үжң¬жЎҶпјҢжүҖжңүж–Үд»¶дјҡиҮӘеҠЁеӨҮд»Ҫ")